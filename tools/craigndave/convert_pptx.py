#!/usr/bin/env python3
"""Convert a Craig 'n' Dave IGCSE PPTX presentation to Markdown.

Usage:
    python tools/craigndave/convert_pptx.py <input.pptx> <slug> [--content]
    # Outputs: markdown/craigndave/<slug>_summary.md
    #          markdown/craigndave/<slug>_content.md  (only with --content)
"""
import re
import sys
from pathlib import Path

from pptx import Presentation

GOING_BEYOND_RE = re.compile(r"going beyond", re.IGNORECASE)
OUT_DIR = Path("markdown/craigndave")
INDENT = ["", "  ", "    ", "      "]


def shape_top(shape) -> int:
    try:
        return shape.top or 0
    except Exception:
        return 0


def slide_texts(slide) -> list[tuple[int, str]]:
    shapes = sorted(
        [s for s in slide.shapes if s.has_text_frame],
        key=shape_top,
    )
    result: list[tuple[int, str]] = []
    for shape in shapes:
        for para in shape.text_frame.paragraphs:
            text = para.text.replace("\x0b", " ").strip()
            if text:
                result.append((para.level, text))
    return result


def is_going_beyond(slide) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame and GOING_BEYOND_RE.search(shape.text_frame.text):
            return True
    return False


def format_slide_md(items: list[tuple[int, str]]) -> str:
    if not items:
        return ""
    lines: list[str] = []
    first = True
    for level, text in items:
        if first and level == 0:
            lines.append(f"### {text}")
            first = False
        else:
            prefix = INDENT[min(level, 3)]
            lines.append(f"{prefix}- {text}")
            first = False
    return "\n".join(lines)


def build_content_md(slides_items: list[list[tuple[int, str]]], title: str, subtitle: str) -> str:
    parts: list[str] = [f"# {title}"]
    if subtitle:
        parts.append(f"*{subtitle}*")
    parts.append("\n---\n")
    for items in slides_items:
        if not items:
            continue
        block = format_slide_md(items)
        if block:
            parts.append(block)
            parts.append("")
    return "\n".join(parts)


def build_summary_md(items: list[tuple[int, str]], title: str, subtitle: str) -> str:
    parts: list[str] = [f"# Summary: {title}"]
    if subtitle:
        parts.append(f"*{subtitle}*")
    parts.append("\n---\n")
    if items:
        parts.append(format_slide_md(items))
    return "\n".join(parts)


def convert(pptx_path: Path, slug: str, with_content: bool = False) -> None:
    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    going_beyond_idx = None
    for i, slide in enumerate(slides):
        if is_going_beyond(slide):
            going_beyond_idx = i
            break

    cutoff = going_beyond_idx if going_beyond_idx is not None else len(slides)
    igcse_slides = slides[:cutoff]

    all_items = [slide_texts(s) for s in igcse_slides]

    title_items = all_items[0] if all_items else []
    title = title_items[0][1] if title_items else slug
    subtitle = title_items[1][1] if len(title_items) > 1 else ""

    content_items = all_items[1:]

    summary_items: list[tuple[int, str]] = []
    for items in reversed(content_items):
        if items:
            summary_items = items
            break

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    if with_content:
        content_path = OUT_DIR / f"{slug}_content.md"
        content_path.write_text(build_content_md(content_items, title, subtitle), encoding="utf-8")
        igcse_count = sum(1 for items in content_items if items)
        print(f"Written: {content_path} ({igcse_count} content slides)")

    summary_path = OUT_DIR / f"{slug}_summary.md"
    summary_path.write_text(build_summary_md(summary_items, title, subtitle), encoding="utf-8")
    print(f"Written: {summary_path}")

    if going_beyond_idx is not None:
        print(f"Excluded: {len(slides) - cutoff} 'Going beyond' slides (from slide {going_beyond_idx + 1})")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <slug> [--content]")
        sys.exit(1)
    with_content = "--content" in sys.argv
    convert(Path(sys.argv[1]), sys.argv[2], with_content)
